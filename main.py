import json
import os
from io import BytesIO
from typing import Literal

from dotenv import load_dotenv
import msal
from openai import OpenAI
from openpyxl import Workbook
from openpyxl.cell.cell import ILLEGAL_CHARACTERS_RE
from openpyxl.styles import Alignment
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pydantic import BaseModel, Field
import requests

load_dotenv()

EXCEL_CELL_CHARACTER_LIMIT = 32767
SLIDE_BREAK = "\n\n--- SLIDE BREAK ---\n\n"
DEFAULT_DATA_ROW_HEIGHT = 15
MINIMUM_COLUMN_WIDTH = 20
GENERATED_BY_AI_SUFFIX = "*"
OUTPUT_DIR = "output"


class PresentationMetadata(BaseModel):
    theme: list[
        Literal[
            "Confidence & Leadership",
            "Financial Literacy",
            "College & Career Prep",
        ]
    ] = Field(
        description=(
            "One or more applicable themes. Only use multiple themes when necessary."
        )
    )
    description: str = Field(
        description="A concise description of the presentation in one or two sentences."
    )
    duration_estimate_minutes: int = Field(
        description=(
            "Estimated total presentation duration in minutes, rounded to the nearest "
            "15 minutes unless over 120 minutes, in which case round to the nearest hour."
        )
    )
    audience: Literal["Middle school", "High school", "College"]
    activity_length_minutes: int = Field(
        description=(
            "Approximate minutes of activity time in the presentation, using the same "
            "rounding rules as duration_estimate_minutes."
        )
    )


def get_openai_client() -> OpenAI:
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise ValueError("OPENAI_API_KEY is missing from the environment")
    return OpenAI(api_key=api_key)


def get_site_id(site_hostname, site_path, headers):
    url = f"https://graph.microsoft.com/v1.0/sites/{site_hostname}:{site_path}"
    resp = requests.get(url, headers=headers)
    resp.raise_for_status()
    site_data = resp.json()
    return site_data["id"]


def get_drive_id(site_id, drive_name, headers):
    url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives"
    resp = requests.get(url, headers=headers)
    resp.raise_for_status()
    drives_data = resp.json()
    for drive in drives_data["value"]:
        if drive["name"] == drive_name:
            return drive["id"]
    raise ValueError("Drive not found")


def get_all_pptx_files(drive_id, headers, item_id="") -> list[str]:
    item_path = f"items/{item_id}" if item_id else "root"
    url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/{item_path}/children"
    resp = requests.get(url, headers=headers)
    resp.raise_for_status()
    items = resp.json()["value"]

    subfolders = [x for x in items if ("folder" in x)]
    subfolder_pptx_files = [
        get_all_pptx_files(drive_id, headers, x["id"]) for x in subfolders
    ]
    pptx_files = [x for x in items if x["name"].lower().endswith(".pptx")]
    return [f for f in pptx_files] + [
        f for subfolder_files in subfolder_pptx_files for f in subfolder_files
    ]


def get_file(drive_id, item_id, headers):
    url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{item_id}"
    resp = requests.get(url, headers=headers)
    resp.raise_for_status()
    return resp.json()


def download_pptx_file_content(drive_id, item_id, headers):
    url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{item_id}/content"
    resp = requests.get(url, headers=headers)
    resp.raise_for_status()
    return resp.content


def extract_slide_text_from_pptx_bytes(pptx_bytes) -> list[str]:
    presentation = Presentation(BytesIO(pptx_bytes))
    slides = []

    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        slides.append("\n".join(slide_text))

    return slides


def generate_ai_metadata(
    client: OpenAI,
    *,
    name: str,
    slide_texts: list[str],
    number_of_slides: int,
    average_words_per_slide: float,
) -> PresentationMetadata:
    entire_slide_text = SLIDE_BREAK.join(slide_texts)

    response = client.responses.parse(
        model="gpt-4.1",
        input=[
            {
                "role": "system",
                "content": (
                    "You classify educational presentation decks. "
                    "Return only schema-valid structured data."
                ),
            },
            {
                "role": "user",
                "content": (
                    f"Presentation name: {name}\n"
                    f"Number of slides: {number_of_slides}\n"
                    f"Average words per slide: {average_words_per_slide:.2f}\n"
                    "Entire slide text:\n"
                    f"{entire_slide_text}"
                ),
            },
        ],
        text_format=PresentationMetadata,
    )

    if response.output_parsed is None:
        raise ValueError(f"OpenAI did not return structured metadata for {name}")

    return response.output_parsed


def sanitize_excel_value(value):
    if not isinstance(value, str):
        return value
    cleaned = ILLEGAL_CHARACTERS_RE.sub("", value)
    return cleaned[:EXCEL_CELL_CHARACTER_LIMIT]


def serialize_object_for_excel(obj: dict) -> dict:
    serialized = {}

    for key, value in obj.items():
        if key == "slide_texts":
            serialized[key] = sanitize_excel_value(SLIDE_BREAK.join(map(str, value)))
        elif isinstance(value, list):
            serialized[key] = sanitize_excel_value(", ".join(map(str, value)))
        elif isinstance(value, (dict, tuple, set)):
            serialized[key] = sanitize_excel_value(json.dumps(value))
        else:
            serialized[key] = sanitize_excel_value(value)

    return serialized


def write_objects_to_excel(
    objects: list[dict], output_filename: str = "workshop_catalog.xlsx"
) -> None:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Presentations"

    output_path = os.path.join(OUTPUT_DIR, output_filename)

    if not objects:
        workbook.save(output_path)
        return

    excel_objects = [serialize_object_for_excel(obj) for obj in objects]
    headers = list(excel_objects[0].keys())
    worksheet.append(headers)

    for column_index, header in enumerate(headers, start=1):
        column_letter = get_column_letter(column_index)
        worksheet.column_dimensions[column_letter].width = max(
            len(str(header)) + 2, MINIMUM_COLUMN_WIDTH
        )

    for obj in excel_objects:
        worksheet.append([obj.get(header) for header in headers])

    for row in worksheet.iter_rows(min_row=2):
        worksheet.row_dimensions[row[0].row].height = DEFAULT_DATA_ROW_HEIGHT
        for cell in row:
            cell.alignment = Alignment(wrap_text=False, vertical="top")

    workbook.save(output_path)


def main():
    TENANT_ID = os.getenv("TENANT_ID")
    CLIENT_ID = os.getenv("CLIENT_ID")
    CLIENT_SECRET = os.getenv("CLIENT_SECRET_VALUE")  # for confidential/server apps
    AUTHORITY = f"https://login.microsoftonline.com/{TENANT_ID}"
    SCOPES = ["https://graph.microsoft.com/.default"]
    SITE_HOSTNAME = os.getenv("SITE_HOSTNAME")
    SITE_PATH = os.getenv("SITE_PATH")
    openai_client = get_openai_client()

    app = msal.ConfidentialClientApplication(
        CLIENT_ID,
        authority=AUTHORITY,
        client_credential=CLIENT_SECRET,
    )

    token = app.acquire_token_for_client(scopes=SCOPES)
    access_token = token["access_token"]

    headers = {"Authorization": f"Bearer {access_token}"}

    SITE_ID = get_site_id(SITE_HOSTNAME, SITE_PATH, headers)
    LIBRARY_DRIVE_ID = get_drive_id(SITE_ID, os.getenv("DRIVE_NAME"), headers)

    # raw_pptx_files = get_all_pptx_files(LIBRARY_DRIVE_ID, headers)
    raw_pptx_files = [
        get_file(LIBRARY_DRIVE_ID, item_id, headers)
        for item_id in [
            "01I7HKCO3RVKMEHQRDR5GZJS6QR56L6LCY",
            # "01I7HKCO4N6ZP6BHCCCVBJDSLM4WQMMQ3Q",
            # "01I7HKCO5IQU3OKDVXEJHYBUP7LAFU4UHH",
            # "01I7HKCO4SPWFCOVNN7JAL4QDULH5PPCYJ",
            # "01I7HKCO6FJRJTI7CXJRAISZFRGAYIPMPU",
            # "01I7HKCO7VKOUI5SISPVGITFBOSTIWTI3H",
        ]
    ]

    base_pptx_slide_data = [
        {
            "id": pptx_file["id"],
            "name": pptx_file["name"],
            "web_url": pptx_file["webUrl"],
            "slide_texts": extract_slide_text_from_pptx_bytes(
                download_pptx_file_content(LIBRARY_DRIVE_ID, pptx_file["id"], headers)
            ),
            "last_modified": pptx_file["lastModifiedDateTime"],
        }
        for pptx_file in raw_pptx_files
    ]

    final_pptx_objects = []
    for pptx_file in base_pptx_slide_data:
        number_of_slides = len(pptx_file["slide_texts"])
        average_words_per_slide = (
            sum(len(slide.split()) for slide in pptx_file["slide_texts"])
            / number_of_slides
        )
        ai_metadata = generate_ai_metadata(
            openai_client,
            name=pptx_file["name"],
            slide_texts=pptx_file["slide_texts"],
            number_of_slides=number_of_slides,
            average_words_per_slide=average_words_per_slide,
        )

        final_pptx_objects.append(
            {
                **pptx_file,
                "number_of_slides": number_of_slides,
                "average_words_per_slide": average_words_per_slide,
                f"theme{GENERATED_BY_AI_SUFFIX}": ai_metadata.theme,
                f"description{GENERATED_BY_AI_SUFFIX}": ai_metadata.description,
                f"duration_estimate_mins{GENERATED_BY_AI_SUFFIX}": ai_metadata.duration_estimate_minutes,
                f"audience{GENERATED_BY_AI_SUFFIX}": ai_metadata.audience,
                f"activity_length_mins{GENERATED_BY_AI_SUFFIX}": ai_metadata.activity_length_minutes,
            }
        )

    write_objects_to_excel(final_pptx_objects)


if __name__ == "__main__":
    main()
