"""Utilities for reading PowerPoint files.

The rest of the app works with plain strings, so this module converts raw PPTX
bytes into a simple `list[str]` where each element is the text from one slide.
"""

from io import BytesIO

from pptx import Presentation


def extract_slide_text_from_pptx_bytes(pptx_bytes) -> list[str]:
    """Extract visible text from each slide in a PowerPoint file.

    Args:
        pptx_bytes: Raw `.pptx` file contents downloaded from Microsoft Graph.

    Returns:
        A list where each item contains the text from one slide. Shapes without
        a usable `.text` attribute are ignored.
    """
    presentation = Presentation(BytesIO(pptx_bytes))
    slides = []

    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            # Many PowerPoint shapes do not contain text. `hasattr(...)` lets us
            # skip image-only and decorative shapes safely.
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        slides.append("\n".join(slide_text))

    return slides
